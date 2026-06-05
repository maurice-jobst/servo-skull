"""Phase 1: Deterministic document extraction using MarkItDown."""
import logging
from pathlib import Path
from typing import Any, Literal

from lexmechanic import Lexmechanic

from servo_skull._utils import setup_logging, _load_config
from servo_skull.llm_providers import expand_env_vars
from servo_skull.models import DocumentExtract

logger = setup_logging(__name__)


class Extractor:
    """Handles deterministic extraction from various file types by wrapping Lexmechanic."""

    SUPPORTED_TYPES = {
        # PDF
        ".pdf": "pdf",
        # Office
        ".docx": "office",
        ".xlsx": "office",
        ".pptx": "office",
        ".doc": "office",
        ".xls": "office",
        ".html": "office",
        ".htm": "office",
        # Image
        ".png": "image",
        ".jpg": "image",
        ".jpeg": "image",
        ".tiff": "image",
        ".gif": "image",
        # Audio / Video transcription
        ".mp3": "audio",
        ".wav": "audio",
        ".m4a": "audio",
        ".flac": "audio",
        ".mp4": "audio",
        ".mkv": "audio",
        ".avi": "audio",
        ".mov": "audio",
        ".webm": "audio",
        # Plain text and chat transcripts
        ".txt": "text",
        ".md": "text",
        ".rst": "text",
        ".json": "text",
    }

    def __init__(self) -> None:
        """Initialize Lexmechanic with configuration."""
        config = _load_config()
        extraction_cfg = config.get("extraction", {})
        expanded_cfg = expand_env_vars(extraction_cfg)

        image_backend = expanded_cfg.pop("image_backend", "tesseract").lower()
        pdf_backend = expanded_cfg.pop("pdf_backend", "pymupdf4llm").lower()

        # Derive standard Ollama url if vlm_endpoint is not defined
        primary_name = config.get("routing", {}).get("primary", "local_gemma")
        primary_cfg = config.get("providers", {}).get(primary_name, {})
        default_endpoint = primary_cfg.get("base_url", "http://127.0.0.1:11434/v1")
        if default_endpoint.endswith("/v1"):
            default_endpoint = default_endpoint[:-3]
        elif default_endpoint.endswith("/v1/"):
            default_endpoint = default_endpoint[:-4]

        vlm_endpoint = expanded_cfg.pop("vlm_endpoint", default_endpoint)
        vlm_model = expanded_cfg.pop("vlm_model", "qwen2.5-vl")
        vlm_prompt = expanded_cfg.pop("vlm_prompt", None)

        # Create Lexmechanic instance with options
        self.md = Lexmechanic(
            image_backend=image_backend,
            pdf_backend=pdf_backend,
            vlm_endpoint=vlm_endpoint,
            vlm_model=vlm_model,
            vlm_prompt=vlm_prompt,
            **expanded_cfg
        )

    def extract(self, file_path: Path) -> DocumentExtract:
        """Extract text and structure from a file."""
        file_path = Path(file_path)

        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        suffix = file_path.suffix.lower()
        if suffix not in self.SUPPORTED_TYPES:
            raise ValueError(f"Unsupported file type: {suffix}")

        doc_type = self.SUPPORTED_TYPES[suffix]
        logger.info(f"Extracting {doc_type} from {file_path.name}")

        text = ""
        metadata: dict[str, Any] = {}
        warnings: list[str] = []
        confidence = 0.90
        extraction_tool = "markitdown"

        # Establish baseline confidence and tool naming per type
        if doc_type == "text":
            confidence = 0.99
            extraction_tool = "text_reader"
        elif doc_type == "pdf":
            confidence = 0.95
            extraction_tool = "pymupdf4llm"
        elif doc_type == "image":
            confidence = 0.60
            extraction_tool = "tesseract_ocr"
        elif doc_type == "audio":
            confidence = 0.75
            extraction_tool = "whisper"

        try:
            file_size = file_path.stat().st_size
        except OSError:
            file_size = 0
        metadata["file_size_bytes"] = file_size

        try:
            # Conversion routed through MarkItDown and our local converters
            result = self.md.convert(str(file_path))
            text = result.text_content
            
            # Check if inline fallback to Tesseract occurred (use explicit 'is True' to avoid MagicMock truthiness)
            if getattr(result, "vlm_fallback", False) is True:
                warnings.append("Ollama VLM failed; fell back to local Tesseract OCR.")
                confidence = 0.60
                extraction_tool = "tesseract_ocr(vlm_fallback)"
        except Exception as e:
            logger.warning(
                f"MarkItDown extraction failed for {file_path.name}: {e}. Trying fallback."
            )
            warnings.append(f"MarkItDown failed: {e}")
            extraction_tool = f"office_parser_fallback({suffix})"

            # Fallback to legacy custom parsers (retained for resilience)
            if suffix == ".docx":
                try:
                    from docx import Document
                    doc = Document(file_path)
                    text = "\n".join(para.text for para in doc.paragraphs)
                    metadata["paragraphs"] = len(doc.paragraphs)
                except Exception as ex:
                    warnings.append(f"DOCX fallback failure: {ex}")
                    confidence = 0.3
            elif suffix in {".xlsx", ".xls"}:
                try:
                    import openpyxl
                    wb = openpyxl.load_workbook(file_path)
                    try:
                        lines = []
                        for sheet in wb.sheetnames:
                            ws = wb[sheet]
                            lines.append(f"Sheet: {sheet}")
                            for row in ws.iter_rows(values_only=True):
                                lines.append(
                                    "\t".join(str(v) if v else "" for v in row)
                                )
                        text = "\n".join(lines)
                        metadata["sheets"] = wb.sheetnames
                    finally:
                        wb.close()
                except Exception as ex:
                    warnings.append(f"XLSX fallback failure: {ex}")
                    confidence = 0.3
            elif suffix == ".pptx":
                try:
                    from pptx import Presentation
                    prs = Presentation(file_path)
                    lines = []
                    for slide_num, slide in enumerate(prs.slides, 1):
                        lines.append(f"Slide {slide_num}:")
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                lines.append(shape.text)
                    text = "\n".join(lines)
                    metadata["slides"] = len(prs.slides)
                except Exception as ex:
                    warnings.append(f"PPTX fallback failure: {ex}")
                    confidence = 0.3

        if not text:
            confidence = 0.0

        return DocumentExtract(
            original_filename=file_path.name,
            document_type=doc_type,
            extracted_text=text,
            metadata=metadata,
            confidence=confidence,
            warnings=warnings,
            extraction_tool=extraction_tool,
        )
